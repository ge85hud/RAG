from pathlib import Path
from typing import List, Dict
import re

from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter


def extract_slide_text(ppt_path: Path) -> List[Dict]:
    """
    Extract text slide-by-slide, keeping slide number and title.
    """
    prs = Presentation(str(ppt_path))
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        title = ""
        content_lines = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            # First text frame often contains title
            if not title:
                title = text
            else:
                content_lines.append(text)

        full_text = "\n".join(content_lines).strip()

        if full_text:
            slides_data.append({
                "slide_number": idx + 1,
                "title": title,
                "text": full_text
            })

    return slides_data


def clean_slide_text(text: str) -> str:
    """
    Normalize slide text (remove excessive bullets and spacing).
    """
    text = re.sub(r"•|\u2022", "-", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def chunk_slide_text(
    slide_text: str,
    chunk_size: int = 400,
    chunk_overlap: int = 80
) -> List[str]:
    """
    Smaller chunks than papers because slides are concise.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n", " ", ""]
    )
    return splitter.split_text(slide_text)


def chunk_lecture_slides(ppt_path: Path) -> List[Dict]:
    """
    Chunk lecture slides into semantically meaningful pieces.
    """
    slides = extract_slide_text(ppt_path)
    documents = []
    chunk_id = 0

    for slide in slides:
        cleaned_text = clean_slide_text(slide["text"])
        chunks = chunk_slide_text(cleaned_text)

        for chunk in chunks:
            documents.append({
                "content": chunk,
                "metadata": {
                    "source": ppt_path.name,
                    "slide_number": slide["slide_number"],
                    "slide_title": slide["title"],
                    "chunk_id": chunk_id,
                    "type": "slide"
                }
            })
            chunk_id += 1

    return documents
